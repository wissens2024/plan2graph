#!/usr/bin/env python3
"""make_ppt.py — 논문 중간 발표자료(.pptx) · 프리미엄 디자인(잉크 네이비 + 골드).

실행: python scripts/make_ppt.py [out.pptx]
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 프리미엄 팔레트 ──
INK = RGBColor(0x10, 0x14, 0x1C)      # 딥 잉크 네이비(표지/섹션)
INK2 = RGBColor(0x1A, 0x20, 0x2C)     # 패널 다크
GOLD = RGBColor(0xC2, 0xA1, 0x4D)     # 골드 액센트
GOLDL = RGBColor(0xE6, 0xD6, 0xA6)    # 라이트 골드
IVORY = RGBColor(0xFB, 0xFA, 0xF6)    # 웜 화이트 배경
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x21, 0x27, 0x33)    # 본문 진회
MUTE = RGBColor(0x8C, 0x91, 0x9C)     # 뮤트 그레이
LINE = RGBColor(0xE7, 0xE3, 0xD8)     # 헤어라인
GREEN = RGBColor(0x3F, 0x6F, 0x52)
RED = RGBColor(0x9C, 0x40, 0x40)
BLUE = RGBColor(0x3B, 0x53, 0x7C)

prs = Presentation()
W, H = Inches(13.333), Inches(7.5)
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
F = "맑은 고딕"
BRAND = "plan2graph · 한국형 소버린 평면도 생성 AI"


def _font(run, size, color, bold=False, italic=False, font=F, spacing=None):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if spacing is not None:                       # 자간(트래킹) — 프리미엄 느낌
        run.font._rPr.set("spc", str(int(spacing * 100)))


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, round_=False, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = el.makeelement(qn('a:outerShdw'),
                            {'blurRad': '90000', 'dist': '38000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val': '0B0E14'})
        al = el.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(al); sh.append(clr); ef.append(sh); el.append(ef)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    return shp


def hline(s, x, y, w, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, x, y, x + w, y)
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=5):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, r in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        txt, kw = (r if isinstance(r, tuple) else (r, {}))
        p.level = kw.get("level", 0)
        if kw.get("bullet"):
            run0 = p.add_run(); _font(run0, kw.get("size", 15), GOLD, bold=True)
            run0.text = "—  "
        run = p.add_run(); run.text = txt
        _font(run, kw.get("size", 15), kw.get("color", SLATE), kw.get("bold", False),
              kw.get("italic", False), kw.get("font", F), kw.get("spc"))
    return tb


def newslide(bg=IVORY):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=bg)
    return s


def header(s, title, kicker=""):
    rect(s, Inches(0.6), Inches(0.62), Inches(0.045), Inches(0.62), fill=GOLD)
    if kicker:
        text(s, Inches(0.82), Inches(0.55), Inches(11), Inches(0.3),
             [(kicker.upper(), {"size": 11, "color": GOLD, "bold": True, "spc": 2.2})])
    text(s, Inches(0.8), Inches(0.78), Inches(11.6), Inches(0.6),
         [(title, {"size": 24, "color": SLATE, "bold": True})])
    hline(s, Inches(0.6), Inches(1.52), Inches(12.13), color=LINE, weight=1.0)
    # 푸터
    hline(s, Inches(0.6), Inches(7.0), Inches(12.13), color=LINE, weight=0.75)
    text(s, Inches(0.6), Inches(7.05), Inches(9), Inches(0.3),
         [(BRAND, {"size": 8.5, "color": MUTE, "spc": 0.6})])
    text(s, Inches(11.6), Inches(7.05), Inches(1.13), Inches(0.3),
         [(f"{len(prs.slides._sldIdLst):02d}", {"size": 8.5, "color": GOLD, "bold": True})],
         align=PP_ALIGN.RIGHT)


def stat(s, x, y, w, label, value, sub=None, accent=GOLD):
    h = Inches(1.45)
    rect(s, x, y, w, h, fill=PAPER, line=LINE, lw=1.0, round_=False, shadow=True)
    rect(s, x, y, w, Inches(0.05), fill=accent)
    text(s, x + Inches(0.28), y + Inches(0.2), w - Inches(0.5), Inches(0.35),
         [(label.upper(), {"size": 10.5, "color": MUTE, "bold": True, "spc": 1.4})])
    text(s, x + Inches(0.26), y + Inches(0.52), w - Inches(0.4), Inches(0.6),
         [(value, {"size": 27, "color": SLATE, "bold": True})])
    if sub:
        text(s, x + Inches(0.28), y + Inches(1.07), w - Inches(0.5), Inches(0.3),
             [(sub, {"size": 10.5, "color": accent, "bold": True})])


def table(s, x, y, w, rows, colw, fs=12.5):
    nr, nc = len(rows), len(rows[0])
    rh = Inches(0.46)
    tbl = s.shapes.add_table(nr, nc, x, y, w, rh * nr).table
    tbl.first_row = False; tbl.horz_banding = False
    for j, cw in enumerate(colw):
        tbl.columns[j].width = cw
    for i, row in enumerate(rows):
        tbl.rows[i].height = rh
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.margin_top = c.margin_bottom = Pt(3)
            c.margin_left = c.margin_right = Pt(9)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = INK if i == 0 else (PAPER if i % 2 else IVORY)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            _font(run, fs, GOLDL if i == 0 else SLATE, bold=(i == 0 or j == 0 and i > 0))
    return tbl


def flow(s, x, y, steps, we=Inches(1.66)):
    cx = x
    for i, st in enumerate(steps):
        b = rect(s, cx, y, we, Inches(0.92), fill=PAPER, line=LINE, lw=1.0, round_=True, shadow=True)
        rect(s, cx, y, we, Inches(0.05), fill=GOLD)
        tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for k, ln in enumerate(st.split("\n")):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = ln
            _font(run, 11, SLATE if k == 0 else MUTE, bold=(k == 0))
        cx += we
        if i < len(steps) - 1:
            t = s.shapes.add_textbox(cx, y, Inches(0.3), Inches(0.92))
            t.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = t.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = "›"; _font(run, 20, GOLD, bold=True)
            cx += Inches(0.3)


def img(s, path, x, y, w=None, h=None):
    return s.shapes.add_picture(path, x, y, width=w, height=h)


# ════════════════════════════════════════════════════════════════════════════
# 1. 표지
s = newslide(INK)
rect(s, Inches(0.45), Inches(0.45), W - Inches(0.9), H - Inches(0.9), fill=None, line=GOLD, lw=1.0)
rect(s, Inches(0.6), Inches(0.6), W - Inches(1.2), H - Inches(1.2), fill=None, line=INK2, lw=0.75)
text(s, Inches(1.1), Inches(0.95), Inches(11), Inches(0.4),
     [("PLAN2GRAPH · 논문 중간 발표", {"size": 12, "color": GOLD, "bold": True, "spc": 3})])
hline(s, Inches(1.13), Inches(2.55), Inches(2.2), color=GOLD, weight=1.5)
text(s, Inches(1.1), Inches(2.75), Inches(11.2), Inches(2),
     [("한국형 소버린", {"size": 46, "color": PAPER, "bold": True}),
      ("아파트 평면도 생성 AI", {"size": 46, "color": GOLDL, "bold": True})], space=4)
text(s, Inches(1.13), Inches(5.0), Inches(11), Inches(0.6),
     [("법규-인식 · 완전 도면(문·창·기구·치수) · 한국 데이터 주권",
       {"size": 16, "color": MUTE, "spc": 0.5})])
text(s, Inches(1.13), Inches(6.35), Inches(11), Inches(0.4),
     [("데이터셋 구축 · 파이프라인 · AI 모델 · 성능        |        2026", {"size": 11, "color": GOLD, "spc": 1})])

# ════════════════════════════════════════════════════════════════════════════
# 2. 연구 배경·목표
s = newslide()
header(s, "연구 배경 · 목표", "Background")
text(s, Inches(0.8), Inches(1.78), Inches(11.7), Inches(1),
     [("건축가가 실제로 쓸 수 있는 ‘완전한’ 한국 아파트 평면도를 생성하는 AI",
       {"size": 19, "bold": True}),
      ("기존 SOTA(RPLAN 류)는 벡터 방·벽 레이아웃까지 — 문·창·기구·치수·법규는 다루지 않는다.",
       {"size": 13.5, "color": MUTE})], space=8)
stat(s, Inches(0.8), Inches(3.25), Inches(3.75), "산출물 ①", "생성형 도면", "화면용 이미지", GOLD)
stat(s, Inches(4.79), Inches(3.25), Inches(3.75), "산출물 ②", "AutoCAD DXF", "작도용 파일", BLUE)
stat(s, Inches(8.78), Inches(3.25), Inches(3.75), "방식", "self-correct", "생성→검증→보정", GREEN)
text(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.8),
     [("핵심 차별성 — 논문 기여", {"size": 14, "bold": True, "color": GOLD, "spc": 1}),
      ("법규-인식 생성: 한국 건축법규(채광 §17①·환기·피난·면적)를 생성에 강제 — 필드 최초", {"size": 14, "bullet": True}),
      ("완전 도면: 레이아웃을 넘어 문·창·기구·치수까지 채운 작도 가능 도면", {"size": 14, "bullet": True}),
      ("한국 소버린: 한국 아파트 데이터로 학습한 자체 소유 엔진(외부 모델 단순 호출 아님)", {"size": 14, "bullet": True})], space=7)

# ════════════════════════════════════════════════════════════════════════════
# 3. 전체 파이프라인
s = newslide()
header(s, "전체 파이프라인", "Overview")
flow(s, Inches(0.62), Inches(1.95),
     ["데이터\n3개 출처", "통일 그래프\ng-0.3", "위상 모델\n방 연결", "기하 모델\nT / G", "neuro-symbolic\n보정·법규", "렌더\n이미지+DXF"])
text(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(0.4),
     [("T · G = ‘같은 파이프라인, 기하모델만 다른’ A/B 비교 (ADR-0007)", {"size": 14, "color": GOLD, "bold": True})])
table(s, Inches(0.8), Inches(3.85), Inches(11.7),
      [["축", "T-라인 (baseline)", "G-라인 (소버린)"],
       ["기하모델", "약한 박스 — ‘위상+약한기하=못그림’ 증명", "강한 diffusion 엔진 (자체)"],
       ["데이터", "자동 추출 그래프", "+ 사람(알바) SVG 보정"],
       ["역할", "비교 기준선", "본 제품"]],
      [Inches(2.0), Inches(5.4), Inches(4.3)])
text(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.4),
     [("비교 지표: FID · 법규준수율 · 완성도 — 기하모델 효과를 격리해 정량 비교", {"size": 12, "color": MUTE})])

# ════════════════════════════════════════════════════════════════════════════
# 4. 핵심 숫자
s = newslide()
header(s, "핵심 숫자 한눈에", "At a Glance")
cw, gap, x0 = Inches(3.78), Inches(0.22), Inches(0.8)
stat(s, x0, Inches(1.85), cw, "받은 원본 · 3출처", "129,007", "AI-Hub·RPLAN·CubiCasa", GOLD)
stat(s, x0 + cw + gap, Inches(1.85), cw, "사용 가능", "94,320", "처분: 사용", GREEN)
stat(s, x0 + 2 * (cw + gap), Inches(1.85), cw, "온전 · AI-Hub APT", "24,706", "64.7% 품질 게이트", BLUE)
stat(s, x0, Inches(3.55), cw, "엔진 · 역할/방", "13 / 18", "한국형 확장", GOLD)
stat(s, x0 + cw + gap, Inches(3.55), cw, "DiffPlanner RPLAN", "FID 1.23", "베이스라인 검증", GREEN)
stat(s, x0 + 2 * (cw + gap), Inches(3.55), cw, "통일 그래프 · G", "40,495", "세대", BLUE)
hline(s, Inches(0.8), Inches(5.5), Inches(11.7), color=GOLD, weight=1.2)
text(s, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.2),
     [("데이터 129,007 → 통일 그래프 40,495세대 → 자체 엔진(13역할/18방) → 완전 도면 + DXF",
       {"size": 16, "bold": True}),
      ("박스 회귀 사망(정답 99% 겹침) → diffusion 엔진 전환 · neuro-symbolic으로 완전성 확보",
       {"size": 13, "color": MUTE})], space=9)

# ════════════════════════════════════════════════════════════════════════════
# 5. 데이터 ① 출처·규모
s = newslide()
header(s, "데이터셋 구축 ① — 출처 · 규모", "Dataset")
text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.4),
     [("3개 공개 데이터셋을 하나의 통일 그래프 규격으로 흡수", {"size": 16, "bold": True})])
table(s, Inches(0.8), Inches(2.4), Inches(11.7),
      [["출처", "받은 원본", "성격", "용도"],
       ["AI-Hub (한국)", "43,219", "한국 아파트·단독·연립 평면도", "소버린 파인튜닝(핵심)"],
       ["RPLAN", "80,788", "대규모 벡터 평면도", "글로벌 사전학습"],
       ["CubiCasa5k", "5,000", "실측 벡터 평면도", "글로벌 사전학습"],
       ["합계", "129,007", "—", "—"]],
      [Inches(2.9), Inches(2.2), Inches(4.3), Inches(2.3)])
stat(s, Inches(0.8), Inches(5.25), Inches(3.75), "받은 원본 합", "129,007", None, GOLD)
stat(s, Inches(4.79), Inches(5.25), Inches(3.75), "사용 가능", "94,320", None, GREEN)
stat(s, Inches(8.78), Inches(5.25), Inches(3.75), "보정·복구 필요", "11,793", None, RED)

# ════════════════════════════════════════════════════════════════════════════
# 6. 데이터 ② 회계
s = newslide()
header(s, "데이터셋 구축 ② — 회계(처분)", "Dataset")
text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.4),
     [("검수 본질 = 숫자 × 카테고리(사용/보정필요/제외) × 사유 · 단일 소스", {"size": 16, "bold": True})])
table(s, Inches(0.8), Inches(2.4), Inches(11.7),
      [["출처", "사용", "보정필요", "제외", "합 = 다운로드"],
       ["AI-Hub", "10,921", "9,412", "22,886", "43,219"],
       ["RPLAN", "80,371", "417", "0", "80,788"],
       ["CubiCasa5k", "3,028", "1,964", "8", "5,000"],
       ["합계", "94,320", "11,793", "22,894", "129,007"]],
      [Inches(2.9), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)])
text(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.6),
     [("원칙", {"size": 14, "bold": True, "color": GOLD, "spc": 1}),
      ("‘보정필요’ = 버리는 게 아니라 살릴 수 있는 데이터 — 자동 재변환 + 사람 보정으로 사용 승급", {"size": 14, "bullet": True}),
      ("AI-Hub는 도면 1장 = 여러 세대 → 도면수 ≠ 세대수(세대 = 실제 추출된 그래프 수)", {"size": 14, "bullet": True})], space=7)

# ════════════════════════════════════════════════════════════════════════════
# 7. 데이터 ③ AI-Hub 처리
s = newslide()
header(s, "데이터셋 구축 ③ — AI-Hub 처리 파이프라인", "Dataset")
flow(s, Inches(0.62), Inches(1.9),
     ["원본 PNG\n라벨", "V2V / YOLO\n검출", "SVG 변환\n완전기하", "사람 보정\n알바", "build\n통일그래프", "검증\n위상·역할"])
text(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.4),
     [("provenance(추출 출처)로 데이터 특성·신뢰도 구분 — 선언적 필터로 구성 분리", {"size": 14, "color": GOLD, "bold": True})])
table(s, Inches(0.8), Inches(3.8), Inches(11.7),
      [["provenance", "의미", "처분"],
       ["dual", "방(SPA)+구조(STR) 둘 다 라벨 — 직접 변환", "사용 (최고품질)"],
       ["spa_only / str_only", "한쪽만 → V2V(YOLO)로 나머지 복구", "사용 (복구)"],
       ["objocr", "OBJ/OCR만(공간 라벨 없음) — 이미지 직접 검출", "보정필요"]],
      [Inches(2.9), Inches(6.2), Inches(2.6)])
text(s, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.8),
     [("V2V(vector→vector) YOLO: SPA mAP 0.90+ · 게이트 통과율로 평가 · G-라인 빌드 18,503 도면 / 40,495 세대",
       {"size": 12, "color": MUTE})])

# ════════════════════════════════════════════════════════════════════════════
# 8. 데이터 ④ 품질 게이트
s = newslide()
header(s, "데이터셋 구축 ④ — 품질 게이트", "Dataset")
text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.4),
     [("‘온전한 데이터만 학습’ — 구성 기반 자동 분류기(좌표 불필요, 검수·학습 단일 소스)", {"size": 16, "bold": True})])
stat(s, Inches(0.8), Inches(2.45), Inches(5.75), "온전 (사용) · AI-Hub APT 38,203 중", "24,706", "64.7%", GREEN)
stat(s, Inches(6.78), Inches(2.45), Inches(5.75), "보정필요", "13,485", "35.3%", RED)
text(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.6),
     [("보정필요 판정 규칙 · 사유 분포 (실측)", {"size": 14, "bold": True, "color": GOLD, "spc": 1}),
      ("현관 ≠ 1 (다세대 병합/누락) — 5,142", {"size": 14, "bullet": True}),
      ("발코니 과다(>4) — 4,904   ·   거실 ≠ 1 — 4,521", {"size": 14, "bullet": True}),
      ("침실 < 화장실(면적 모순) — 3,060   ·   거실 오라벨 — 634   ·   기타 과다 — 527", {"size": 14, "bullet": True}),
      ("온전 데이터 방 수: 중앙값 14 · p95 17 → 엔진 방 수용량 18로 결정", {"size": 14, "bold": True})], space=8)

# ════════════════════════════════════════════════════════════════════════════
# 9. 데이터 ⑤ 구성(variant)
s = newslide()
header(s, "데이터셋 구축 ⑤ — 구성(variant)이 곧 실험 변수", "Dataset")
text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.4),
     [("같은 엔진을 다른 데이터 구성으로 학습 → 다른 모델 → 도면 품질 차이 = 비교의 본체", {"size": 16, "bold": True})])
table(s, Inches(0.8), Inches(2.4), Inches(11.7),
      [["라인", "구성", "세대수", "보정 종류"],
       ["AI-Hub (T)", "dual", "7,465", "자동(재변환) — T·G 공유"],
       ["AI-Hub (T)", "dual + 보정", "23,679", "자동(재변환)"],
       ["AI-Hub (G)", "dual", "8,700", "자동 + 알바(사람 SVG)"],
       ["AI-Hub (G)", "dual + 보정", "27,683", "자동 + 알바"],
       ["글로벌(사전학습)", "RPLAN / CubiCasa", "56,053 / —", "—"]],
      [Inches(2.9), Inches(3.2), Inches(2.4), Inches(3.2)])
text(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(1),
     [("자동(neuro-symbolic 재변환)은 T·G 공유 · 알바(사람 SVG)는 G에만 추가 · 매트릭스 = 구성 × 기하모델 × 사전학습",
       {"size": 12.5, "color": MUTE})])

# ════════════════════════════════════════════════════════════════════════════
# 10. AI 모델 ① 박스회귀 폐기
s = newslide()
header(s, "AI 모델 ① — 박스 회귀(폐기): 측정으로 사망", "AI Model")
text(s, Inches(0.8), Inches(1.78), Inches(11.7), Inches(0.4),
     [("첫 시도: 위상 그래프 → 방별 박스 좌표 회귀(Transformer)", {"size": 16, "bold": True})])
stat(s, Inches(0.8), Inches(2.5), Inches(3.75), "정답 박스 겹침", "99%", "표현 한계", RED)
stat(s, Inches(4.79), Inches(2.5), Inches(3.75), "데이터 4.6배·손실", "효과 0", "동일 붕괴", RED)
stat(s, Inches(8.78), Inches(2.5), Inches(3.75), "배치 실현율", "62%", "방↑일수록 33%", RED)
text(s, Inches(0.8), Inches(4.35), Inches(11.7), Inches(2.4),
     [("왜 죽었나 — 실측 근거", {"size": 14, "bold": True, "color": GOLD, "spc": 1}),
      ("정답(ground-truth) 박스부터 99% 겹침 — 모델이 아니라 ‘박스 표현’ 자체의 한계", {"size": 14, "bullet": True}),
      ("데이터 4.6배 증량·겹침 억제 손실에도 동일하게 박스 붕괴 · 문 증발·화장실 고립", {"size": 14, "bullet": True}),
      ("교훈: 위상/모델 문제가 아니라 ‘기하 실현’이 병목 → 생성형 기하(diffusion)로 전환", {"size": 14, "bold": True})], space=8)

# ════════════════════════════════════════════════════════════════════════════
# 11. AI 모델 ② 소버린 엔진
s = newslide()
header(s, "AI 모델 ② — 소버린 엔진(SOTA 부품 조립 + 한국형 확장)", "AI Model")
text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.4),
     [("단일 생성 엔진 — SOTA의 ‘기법(부품)’만 인용·조립 + 우리 novelty", {"size": 16, "bold": True})])
table(s, Inches(0.8), Inches(2.4), Inches(11.7),
      [["빌린 부품 (인용 base)", "역할"],
       ["DiffPlanner (2025)", "벡터 직접 diffusion 골격 · 경계조건 (코드 보유·검증)"],
       ["GSDiff (AAAI’25)", "정렬 손실 · 자기지도(연결 포인트)"],
       ["HouseDiffusion (CVPR’23)", "이산 + 연속 디노이즈"],
       ["FMLM (2026)", "유효성 제약(markup)"]],
      [Inches(3.9), Inches(7.8)])
text(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.8),
     [("한국형 확장 — 구현 완료·검증", {"size": 14, "bold": True, "color": GOLD, "spc": 1}),
      ("3-stage(node→adjacency→partitioning): num_category 6→13(한국 역할) · max_rooms 8→18", {"size": 14, "bullet": True}),
      ("RPLAN 사전학습 → 한국 파인튜닝(전이학습) · 학습→샘플→도면+DXF 파이프라인 E2E 검증", {"size": 14, "bullet": True})], space=7)

# ════════════════════════════════════════════════════════════════════════════
# 12. AI 모델 ③ neuro-symbolic
s = newslide()
header(s, "AI 모델 ③ — neuro-symbolic 완성", "AI Model")
text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.45),
     [("엔진은 방 레이아웃까지 — 문·창·치수·기구·법규는 ‘있는 데이터로 의미 추론’해 채움", {"size": 16, "bold": True})])
table(s, Inches(0.8), Inches(2.45), Inches(11.7),
      [["채울 공백", "추론 방법"],
       ["축척(scale)", "표준 여닫이문 폭 800mm 앵커 → mm/px (절대값은 사람 길이입력 보조)"],
       ["문 방향 / 위치", "방 인접 관계에서 유도(맞닿은 경계 중점)"],
       ["창문", "거주방의 외곽 접한 벽에 배치"],
       ["기구(침대·소파·싱크…)", "역할별 규칙 카탈로그(벽에 붙여·문 반대쪽)"],
       ["법규 준수", "rules_legal: 채광 §17① · 환기 · 피난 · 면적 (국가법령 API)"]],
      [Inches(3.4), Inches(8.3)])
text(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
     [("생성 → 검증(rules_legal) → 보정 → 재생성 루프 · 거친 출력은 학습 개선 신호", {"size": 12.5, "color": MUTE})])

# ════════════════════════════════════════════════════════════════════════════
# 13. 완전 도면 예시 (이미지)
s = newslide()
header(s, "완전 도면 예시 — neuro-symbolic으로 채운 결과", "Results")
rect(s, Inches(0.62), Inches(1.7), Inches(7.7), Inches(5.05), fill=PAPER, line=LINE, lw=1.0, shadow=True)
img(s, "docs/figs/complete_drawing.png", Inches(0.78), Inches(1.82), h=Inches(4.8))
text(s, Inches(8.62), Inches(1.78), Inches(4.1), Inches(4),
     [("이 도면 1장에 담긴 것", {"size": 15, "bold": True, "color": GOLD, "spc": 0.8}),
      ("방 + 면적(㎡) · 외벽/내벽", {"size": 13.5, "bullet": True}),
      ("가구 — 침대·소파·TV·식탁·싱크/레인지·냉장고·변기·세면·샤워·신발장", {"size": 13.5, "bullet": True}),
      ("창(파랑) · 문(빨강) · 치수 12,653 × 9,513 mm", {"size": 13.5, "bullet": True}),
      ("척도 — 여닫이문 폭 800mm 앵커 추론(11.9 mm/px)", {"size": 13.5, "bullet": True}),
      ("같은 도면 → AutoCAD DXF 동시 출력", {"size": 13.5, "bullet": True, "bold": True, "color": GREEN})], space=9)
hline(s, Inches(8.62), Inches(5.75), Inches(4.1), color=LINE)
text(s, Inches(8.62), Inches(5.85), Inches(4.1), Inches(1),
     [("SOTA는 방·벽 레이아웃까지 — 우리는 가구·치수·DXF까지 채운 ‘완전 도면’",
       {"size": 12, "color": MUTE, "italic": True})])

# ════════════════════════════════════════════════════════════════════════════
# 14. 성능·측정
s = newslide()
header(s, "성능 · 측정", "Results")
stat(s, Inches(0.8), Inches(1.8), Inches(3.75), "DiffPlanner RPLAN", "FID 1.23", "pytorch-fid", GREEN)
stat(s, Inches(4.79), Inches(1.8), Inches(3.75), "clean-fid", "1.34", "동일 셋", GREEN)
stat(s, Inches(8.78), Inches(1.8), Inches(3.75), "평가셋", "12,002", "RPLAN test", BLUE)
text(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.2),
     [("측정으로 확정된 사실 — 흔들지 않는 베이스라인", {"size": 14, "bold": True, "color": GOLD, "spc": 1}),
      ("DiffPlanner 베이스라인 재현·검증: RPLAN FID 1.23(pytorch-fid) / 1.34(clean-fid), 12,002 test", {"size": 14, "bullet": True}),
      ("균형 벤치마크(주거형태 매크로): 신경망 0.188 > 규칙기반 0.205 — ‘규칙기반 최강’은 APT 편중 착시", {"size": 14, "bullet": True}),
      ("V2V 확장은 생성 학습엔 역효과: adj_L1 0.074→0.104 · 법규 0.94→0.37 (수량 ≠ 품질)", {"size": 14, "bullet": True}),
      ("SOTA 4편 전부 벡터 방·벽 레이아웃까지 — 창·가구·치수는 아무도 안 함(최신 FMLM 포함)", {"size": 14, "bullet": True}),
      ("평가 축 = FID(품질) + 법규준수율(SOTA≈0) + 완성도", {"size": 14, "bold": True})], space=9)

# ════════════════════════════════════════════════════════════════════════════
# 15. 시스템·도구
s = newslide()
header(s, "시스템 · 도구 — 재현 가능한 프로그램", "System")
text(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.4),
     [("결과물 = ‘내가 돌린 숫자’가 아니라 ‘비교 가능한 프로그램’ (단일 소스 · A/B 한 화면)", {"size": 16, "bold": True})])
table(s, Inches(0.8), Inches(2.45), Inches(11.7),
      [["구성요소", "내용"],
       ["Streamlit 대시보드", "검수·회계·데이터구성·학습/재학습 조작·도면생성·법규 DB"],
       ["도면 생성 화면", "사전학습 × 파인튜닝(데이터 구성+세대수) → 학습 → 샘플 → 이미지+DXF"],
       ["웹 보정 에디터", "문 방향 클릭·벽 드래그 즉시반영 — Streamlit 한계 보완"],
       ["cadrender 공용 코어", "T·G 공용 렌더 → 이미지(matplotlib) + DXF(ezdxf) + 자기교정"]],
      [Inches(3.4), Inches(8.3)])
text(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.4),
     [("인프라: 13/18 아키텍처 패치·검증 · 회계 영속 캐시(27s → 0.05s)", {"size": 13, "bullet": True}),
      ("하이브리드 프론트: 분석·비교 = Streamlit · 편집(클릭·드래그) = 웹 SVG", {"size": 13, "bullet": True})], space=7)

# ════════════════════════════════════════════════════════════════════════════
# 16. 현재 상태
s = newslide()
header(s, "현재 상태 — 학습 진행 중", "Status")
stat(s, Inches(0.8), Inches(1.85), Inches(3.75), "엔진 아키텍처", "완료", "검증됨", GREEN)
stat(s, Inches(4.79), Inches(1.85), Inches(3.75), "샘플→렌더(+DXF)", "완료", "E2E 검증", GREEN)
stat(s, Inches(8.78), Inches(1.85), Inches(3.75), "엔진 학습", "진행 중", "RPLAN→한국", GOLD)
text(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(3),
     [("진척", {"size": 14, "bold": True, "color": GOLD, "spc": 1}),
      ("품질 게이트·통일 그래프·엔진 데이터 변환(온전·13/18)·아키텍처 — 완료", {"size": 14, "bullet": True}),
      ("3-stage 학습 게이트 통과(한국 데이터, loss 수렴·무오류) · 샘플→렌더 파이프라인 완성", {"size": 14, "bullet": True}),
      ("현재: RPLAN 사전학습 + 한국 파인튜닝 진행 · 짧은 런으로 첫 도면 선행 확보 중", {"size": 14, "bullet": True}),
      ("렌더 검증(정답 그래프 기준): 방·문·창·기구·치수·DXF 전부 생성 — 파이프라인 작동 확인", {"size": 14, "bullet": True})], space=8)

# ════════════════════════════════════════════════════════════════════════════
# 17. 향후 계획·논문
s = newslide()
header(s, "향후 계획 · 논문", "Next")
text(s, Inches(0.8), Inches(1.9), Inches(5.9), Inches(4.5),
     [("로드맵", {"size": 15, "bold": True, "color": GOLD, "spc": 1}),
      ("학습 완료 → 동결 test 샘플링", {"size": 14, "bullet": True}),
      ("T ∥ G · 데이터구성 · 사전학습 비교", {"size": 14, "bullet": True}),
      ("FID · 법규준수율 · 완성도 정량 평가", {"size": 14, "bullet": True}),
      ("알바 보정 → 더 좋은 그래프 → 재학습(반복 루프)", {"size": 14, "bullet": True}),
      ("성공 기준 = FID 동급 AND 법규·완전 압도", {"size": 14, "bold": True})], space=11)
rect(s, Inches(6.95), Inches(1.85), Inches(5.78), Inches(4.5), fill=INK, round_=True, shadow=True)
rect(s, Inches(6.95), Inches(1.85), Inches(5.78), Inches(0.06), fill=GOLD)
text(s, Inches(7.3), Inches(2.15), Inches(5.2), Inches(4),
     [("목표 논문", {"size": 16, "bold": True, "color": GOLDL, "spc": 1}),
      ("A · 한국 아파트 데이터셋 · 벤치마크", {"size": 14, "bullet": True, "color": PAPER}),
      ("B · 완전 도면 생성(문·창·기구·치수) 방법", {"size": 14, "bullet": True, "color": PAPER}),
      ("C · 법규-인식 생성 (필드 최초)", {"size": 14, "bullet": True, "color": PAPER}),
      ("D · 자체 소버린 엔진(SOTA 부품 + novelty)", {"size": 14, "bullet": True, "color": PAPER})], space=13)

# ════════════════════════════════════════════════════════════════════════════
# 18. 마무리
s = newslide(INK)
rect(s, Inches(0.45), Inches(0.45), W - Inches(0.9), H - Inches(0.9), fill=None, line=GOLD, lw=1.0)
hline(s, Inches(1.13), Inches(2.7), Inches(2.2), color=GOLD, weight=1.5)
text(s, Inches(1.1), Inches(2.9), Inches(11.2), Inches(1.6),
     [("요약", {"size": 13, "color": GOLD, "bold": True, "spc": 3}),
      ("한국 데이터 129,007 → 통일 그래프 → 자체 소버린 엔진(13/18) → 완전 도면 + DXF",
       {"size": 23, "color": PAPER, "bold": True})], space=12)
text(s, Inches(1.13), Inches(4.7), Inches(11), Inches(1.5),
     [("박스 회귀 사망 → diffusion 엔진 전환 · neuro-symbolic으로 SOTA가 안 하는 완전성 확보",
       {"size": 14, "color": GOLDL}),
      ("차별성 = 법규-인식 · 완전 도면 · 한국 소버린 — 비교 가능한 프로그램으로 검증 진행 중",
       {"size": 14, "color": MUTE})], space=10)

out = sys.argv[1] if len(sys.argv) > 1 else "docs/plan2graph_중간발표.pptx"
prs.save(out)
print("saved:", out, "·", len(prs.slides._sldIdLst), "slides")
